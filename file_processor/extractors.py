from __future__ import annotations

import csv
import json
import xml.etree.ElementTree as ET
from datetime import date, datetime, time
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Callable


Extractor = Callable[[Path], dict[str, Any]]


class MissingDependencyError(RuntimeError):
    """Raised when an optional parser package is needed but not installed."""


class _HTMLTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self._parts.append(text)

    @property
    def text(self) -> str:
        return "\n".join(self._parts)


def supported_extensions() -> list[str]:
    return sorted(EXTRACTORS)


def extract(path: Path) -> dict[str, Any]:
    extractor = EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise ValueError(f"Unsupported file type: {path.suffix or '<none>'}")
    return extractor(path)


def extract_text(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return {"kind": "text", "text": text}


def extract_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as file:
        data = json.load(file)
    return {"kind": "json", "data": data}


def extract_csv(path: Path, delimiter: str = ",") -> dict[str, Any]:
    with path.open("r", encoding="utf-8-sig", newline="") as file:
        sample = file.read(4096)
        file.seek(0)
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|") if sample.strip() else None
        reader = csv.DictReader(file, dialect=dialect) if dialect else csv.DictReader(file, delimiter=delimiter)
        rows = list(reader)
    return {"kind": "table", "rows": rows, "row_count": len(rows)}


def extract_xml(path: Path) -> dict[str, Any]:
    tree = ET.parse(path)
    root = tree.getroot()
    return {
        "kind": "xml",
        "root_tag": root.tag,
        "text": " ".join(part.strip() for part in root.itertext() if part.strip()),
    }


def extract_html(path: Path) -> dict[str, Any]:
    parser = _HTMLTextParser()
    parser.feed(path.read_text(encoding="utf-8", errors="replace"))
    return {"kind": "html", "text": parser.text}


def extract_xlsx(path: Path) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise MissingDependencyError("Install openpyxl to process Excel files.") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheets: dict[str, list[list[Any]]] = {}
    for worksheet in workbook.worksheets:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            rows.append([normalize_excel_cell(cell) for cell in row])
        sheets[worksheet.title] = rows
    workbook.close()
    return {"kind": "workbook", "sheets": sheets}


def normalize_excel_cell(value: Any) -> Any:
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return value


def extract_docx(path: Path) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise MissingDependencyError("Install python-docx to process Word files.") from exc

    document = Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    tables = [
        [[cell.text for cell in row.cells] for row in table.rows]
        for table in document.tables
    ]
    return {"kind": "document", "paragraphs": paragraphs, "tables": tables, "text": "\n".join(paragraphs)}


def extract_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise MissingDependencyError("Install python-pptx to process PowerPoint files.") from exc

    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
        slides.append({"slide": index, "text": "\n".join(text_parts)})
    return {"kind": "presentation", "slides": slides}


def extract_pdf(path: Path) -> dict[str, Any]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise MissingDependencyError("Install pypdf to process PDF files.") from exc

    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return {"kind": "pdf", "pages": pages, "text": "\n".join(pages)}


EXTRACTORS: dict[str, Extractor] = {
    ".csv": extract_csv,
    ".docx": extract_docx,
    ".htm": extract_html,
    ".html": extract_html,
    ".json": extract_json,
    ".log": extract_text,
    ".md": extract_text,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".tsv": lambda path: extract_csv(path, delimiter="\t"),
    ".txt": extract_text,
    ".xlsm": extract_xlsx,
    ".xlsx": extract_xlsx,
    ".xml": extract_xml,
}
